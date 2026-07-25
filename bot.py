import os
import logging
import tempfile
import subprocess
import json
import requests
import uuid
import zipfile
import shutil
import qrcode
from datetime import datetime
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import Application, CommandHandler, MessageHandler, CallbackQueryHandler, filters, ContextTypes
import yt_dlp
from pydub import AudioSegment
from PIL import Image
import fitz  # PyMuPDF for PDF handling
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pyzbar.pyzbar import decode
import io

# ========== CONFIGURATION ==========
TOKEN = "8866299232:AAECrRPPu5cfRMxx3J1i4CkICw4F4G861DA"  # Replace with your token

CHANNELS = [
    {"name": "Pykillinux", "url": "https://t.me/pykillinux"},
    {"name": "Music Search", "url": "https://t.me/music_search"}
]

ADMIN_IDS = [310141017]  # Replace with your user ID

# ========== TEMPORARY STORAGE ==========
temp_storage = {}
user_sessions = {}  # Track user session state

# ========== LANGUAGE DICTIONARY ==========
LANG = {
    'fa': {
        'welcome': "🎯 **ربات چندمنظوره**\n\n"
                   "برای استفاده از این ربات، لطفاً ابتدا در کانال‌های زیر عضو شوید:\n\n",
        'join_channels': "• {name}\n",
        'check_btn': "✅ بررسی عضویت",
        'subscribed': "✅ **تبریک! شما در تمام کانال‌ها عضو هستید.**\n\n"
                      "🎯 از منوی زیر یکی از گزینه‌ها را انتخاب کنید:",
        'not_subscribed': "❌ **شما هنوز در تمام کانال‌ها عضو نشده‌اید.**\n\n"
                          "لطفاً روی لینک کانال‌های زیر کلیک کرده و عضو شوید، سپس دوباره روی دکمه **بررسی عضویت** کلیک کنید.",
        'not_allowed': "❌ **دسترسی محدود شده است.**\n\nلطفاً ابتدا در کانال‌های زیر عضو شوید:",
        'help': "برای راهنمایی بیشتر از /start استفاده کنید.",
        'main_menu': "🎯 **منوی اصلی**\n\nلطفاً یکی از گزینه‌های زیر را انتخاب کنید:",
        'cancel': "❌ عملیات لغو شد.",
        'lang_changed': "🌐 زبان به فارسی تغییر کرد.",
        'lang_btn': "🌐 تغییر زبان",
        'processing': "🔄 در حال پردازش...",
        'download_complete': "✅ عملیات با موفقیت انجام شد!",
        'error': "❌ خطا رخ داد. لطفاً دوباره تلاش کنید.",
        
        # Menu items
        'menu_video_convert': "🎬 تبدیل ویدیو به MP4",
        'menu_audio_convert': "🎵 تبدیل صدا به MP3",
        'menu_youtube': "📥 دانلود از یوتیوب",
        'menu_instagram': "📸 دانلود از اینستاگرام",
        'menu_image_convert': "🖼 تبدیل تصویر (JPG/PNG)",
        'menu_image_to_pdf': "📄 تصویر به PDF",
        'menu_pdf_to_image': "📄 PDF به تصویر",
        'menu_qr_generate': "🔲 تولید QR Code",
        'menu_qr_read': "🔍 خواندن QR Code",
        'menu_zip': "📦 فشرده‌سازی فایل‌ها",
        'menu_unzip': "📂 خارج‌سازی فایل‌ها",
        'menu_doc_to_pdf': "📄 تبدیل سند به PDF",
        'menu_text_to_pdf': "📝 متن به PDF",
        'menu_audio_extract': "🎬 استخراج صدا از ویدیو",
        'menu_unit_convert': "📏 تبدیل واحدها",
        'menu_link_download': "🔗 دانلود از لینک",
        'back': "🔙 بازگشت به منوی اصلی",
        
        # Feature-specific messages
        'send_video': "🎬 لطفاً یک ویدیو ارسال کنید.",
        'send_audio': "🎵 لطفاً یک فایل صوتی ارسال کنید.",
        'send_image': "🖼 لطفاً یک تصویر ارسال کنید.",
        'send_images': "🖼 لطفاً چند تصویر ارسال کنید (حداکثر ۵ عدد). بعد از ارسال، روی دکمه **ایجاد PDF** کلیک کنید.",
        'send_pdf': "📄 لطفاً یک فایل PDF ارسال کنید.",
        'send_document': "📄 لطفاً یک فایل سند (Word/PPT/Excel) ارسال کنید.",
        'send_text': "📝 لطفاً متن خود را ارسال کنید.",
        'send_link': "🔗 لطفاً لینک خود را ارسال کنید.",
        'send_youtube': "📥 لطفاً لینک یوتیوب را ارسال کنید.",
        'send_instagram': "📸 لطفاً لینک اینستاگرام را ارسال کنید.",
        'send_qr_text': "🔲 لطفاً متن مورد نظر برای تولید QR Code را ارسال کنید.",
        'send_qr_image': "🔍 لطفاً تصویر حاوی QR Code را ارسال کنید.",
        'send_zip_files': "📦 لطفاً فایل‌های مورد نظر برای فشرده‌سازی را ارسال کنید (حداکثر ۵ عدد). سپس روی دکمه **فشرده‌سازی** کلیک کنید.",
        'send_unzip': "📂 لطفاً فایل ZIP را ارسال کنید.",
        'convert_to': "تبدیل به",
        'choose_quality': "کیفیت را انتخاب کنید:",
        'pdf_created': "✅ PDF با موفقیت ساخته شد!",
        'qr_created': "✅ QR Code با موفقیت ساخته شد!",
        'qr_read_result': "🔍 نتیجه خواندن QR Code:\n\n{text}",
        'zip_created': "✅ فایل ZIP با موفقیت ساخته شد!",
        'unzip_done': "✅ فایل‌ها با موفقیت خارج‌سازی شدند!",
        'unit_result': "📏 نتیجه تبدیل:\n\n{result}",
        'audio_extracted': "🎵 صدای ویدیو با موفقیت استخراج شد!",
        'file_not_found': "❌ فایل پیدا نشد. لطفاً دوباره تلاش کنید.",
        'max_files': "❌ حداکثر ۵ فایل مجاز است.",
    },
    'en': {
        'welcome': "🎯 **Multipurpose Bot**\n\n"
                   "To use this bot, please join the following channels first:\n\n",
        'join_channels': "• {name}\n",
        'check_btn': "✅ Check Subscription",
        'subscribed': "✅ **Congratulations! You are subscribed to all channels.**\n\n"
                      "🎯 Select an option from the menu below:",
        'not_subscribed': "❌ **You haven't joined all channels yet.**\n\n"
                          "Please click the channel links below and join, then click the **Check Subscription** button again.",
        'not_allowed': "❌ **Access restricted.**\n\nPlease join the channels below first:",
        'help': "Use /start for more information.",
        'main_menu': "🎯 **Main Menu**\n\nPlease select one of the options below:",
        'cancel': "❌ Operation cancelled.",
        'lang_changed': "🌐 Language changed to English.",
        'lang_btn': "🌐 Change Language",
        'processing': "🔄 Processing...",
        'download_complete': "✅ Operation completed successfully!",
        'error': "❌ An error occurred. Please try again.",
        
        # Menu items
        'menu_video_convert': "🎬 Convert Video to MP4",
        'menu_audio_convert': "🎵 Convert Audio to MP3",
        'menu_youtube': "📥 Download from YouTube",
        'menu_instagram': "📸 Download from Instagram",
        'menu_image_convert': "🖼 Image Convert (JPG/PNG)",
        'menu_image_to_pdf': "📄 Image to PDF",
        'menu_pdf_to_image': "📄 PDF to Image",
        'menu_qr_generate': "🔲 Generate QR Code",
        'menu_qr_read': "🔍 Read QR Code",
        'menu_zip': "📦 Compress Files (ZIP)",
        'menu_unzip': "📂 Extract Files (ZIP)",
        'menu_doc_to_pdf': "📄 Document to PDF",
        'menu_text_to_pdf': "📝 Text to PDF",
        'menu_audio_extract': "🎬 Extract Audio from Video",
        'menu_unit_convert': "📏 Unit Converter",
        'menu_link_download': "🔗 Download from Link",
        'back': "🔙 Back to Main Menu",
        
        # Feature-specific messages
        'send_video': "🎬 Please send a video.",
        'send_audio': "🎵 Please send an audio file.",
        'send_image': "🖼 Please send an image.",
        'send_images': "🖼 Please send multiple images (max 5). After sending, click **Create PDF** button.",
        'send_pdf': "📄 Please send a PDF file.",
        'send_document': "📄 Please send a document (Word/PPT/Excel).",
        'send_text': "📝 Please send your text.",
        'send_link': "🔗 Please send your link.",
        'send_youtube': "📥 Please send a YouTube link.",
        'send_instagram': "📸 Please send an Instagram link.",
        'send_qr_text': "🔲 Please send the text to generate a QR Code.",
        'send_qr_image': "🔍 Please send an image containing a QR Code.",
        'send_zip_files': "📦 Please send files to compress (max 5). Then click **Compress** button.",
        'send_unzip': "📂 Please send a ZIP file.",
        'convert_to': "Convert to",
        'choose_quality': "Choose quality:",
        'pdf_created': "✅ PDF created successfully!",
        'qr_created': "✅ QR Code created successfully!",
        'qr_read_result': "🔍 QR Code read result:\n\n{text}",
        'zip_created': "✅ ZIP file created successfully!",
        'unzip_done': "✅ Files extracted successfully!",
        'unit_result': "📏 Conversion result:\n\n{result}",
        'audio_extracted': "🎵 Audio extracted successfully!",
        'file_not_found': "❌ File not found. Please try again.",
        'max_files': "❌ Maximum 5 files allowed.",
    }
}

# ========== LOGGING ==========
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# ========== USER DATA ==========
user_lang = {}
user_subscribed = {}
user_sessions = {}
USER_DATA_FILE = 'user_data.json'

def load_user_data():
    if os.path.exists(USER_DATA_FILE):
        with open(USER_DATA_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {}

def save_user_data(data):
    with open(USER_DATA_FILE, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

user_data = load_user_data()

def track_user(user_id, username=None, first_name=None, last_name=None, action=None):
    now = datetime.now().isoformat()
    uid = str(user_id)
    if uid not in user_data:
        user_data[uid] = {
            'first_seen': now,
            'user_id': user_id,
            'username': username,
            'first_name': first_name,
            'last_name': last_name,
            'total_messages': 0,
            'actions': {},
            'last_active': now,
            'subscribed': False,
        }
    user = user_data[uid]
    user['last_active'] = now
    user['total_messages'] += 1
    if username:
        user['username'] = username
    if first_name:
        user['first_name'] = first_name
    if last_name:
        user['last_name'] = last_name
    if action:
        user['actions'][action] = user['actions'].get(action, 0) + 1
    save_user_data(user_data)
    return user

# ========== SUBSCRIPTION CHECK ==========
async def check_subscription(user_id: int, context: ContextTypes.DEFAULT_TYPE) -> bool:
    for channel in CHANNELS:
        try:
            channel_identifier = channel["url"].split("/")[-1]
            member = await context.bot.get_chat_member(
                chat_id=f"@{channel_identifier}",
                user_id=user_id
            )
            if member.status in ["left", "kicked"]:
                return False
        except Exception as e:
            logger.error(f"Error checking {channel['name']}: {e}")
            return False
    return True

# ========== MENU BUILDERS ==========
def build_main_menu(lang: str):
    keyboard = [
        [InlineKeyboardButton(LANG[lang]['menu_video_convert'], callback_data="menu_video_convert")],
        [InlineKeyboardButton(LANG[lang]['menu_audio_convert'], callback_data="menu_audio_convert")],
        [InlineKeyboardButton(LANG[lang]['menu_youtube'], callback_data="menu_youtube")],
        [InlineKeyboardButton(LANG[lang]['menu_instagram'], callback_data="menu_instagram")],
        [InlineKeyboardButton(LANG[lang]['menu_image_convert'], callback_data="menu_image_convert")],
        [InlineKeyboardButton(LANG[lang]['menu_image_to_pdf'], callback_data="menu_image_to_pdf")],
        [InlineKeyboardButton(LANG[lang]['menu_pdf_to_image'], callback_data="menu_pdf_to_image")],
        [InlineKeyboardButton(LANG[lang]['menu_qr_generate'], callback_data="menu_qr_generate")],
        [InlineKeyboardButton(LANG[lang]['menu_qr_read'], callback_data="menu_qr_read")],
        [InlineKeyboardButton(LANG[lang]['menu_zip'], callback_data="menu_zip")],
        [InlineKeyboardButton(LANG[lang]['menu_unzip'], callback_data="menu_unzip")],
        [InlineKeyboardButton(LANG[lang]['menu_doc_to_pdf'], callback_data="menu_doc_to_pdf")],
        [InlineKeyboardButton(LANG[lang]['menu_text_to_pdf'], callback_data="menu_text_to_pdf")],
        [InlineKeyboardButton(LANG[lang]['menu_audio_extract'], callback_data="menu_audio_extract")],
        [InlineKeyboardButton(LANG[lang]['menu_unit_convert'], callback_data="menu_unit_convert")],
        [InlineKeyboardButton(LANG[lang]['menu_link_download'], callback_data="menu_link_download")],
        [InlineKeyboardButton(LANG[lang]['lang_btn'], callback_data="toggle_lang")]
    ]
    return InlineKeyboardMarkup(keyboard)

def build_back_button(lang: str):
    keyboard = [[InlineKeyboardButton(LANG[lang]['back'], callback_data="main_menu")]]
    return InlineKeyboardMarkup(keyboard)

def build_unit_converter_keyboard(lang: str):
    keyboard = [
        [InlineKeyboardButton("📏 Length", callback_data="unit_length")],
        [InlineKeyboardButton("⚖️ Weight", callback_data="unit_weight")],
        [InlineKeyboardButton("🌡️ Temperature", callback_data="unit_temperature")],
        [InlineKeyboardButton("📐 Area", callback_data="unit_area")],
        [InlineKeyboardButton("🧊 Volume", callback_data="unit_volume")],
        [InlineKeyboardButton(LANG[lang]['back'], callback_data="main_menu")]
    ]
    return InlineKeyboardMarkup(keyboard)

def build_channel_keyboard(lang: str):
    keyboard = []
    for channel in CHANNELS:
        keyboard.append([InlineKeyboardButton(f"📢 {channel['name']}", url=channel["url"])])
    keyboard.append([InlineKeyboardButton(LANG[lang]['check_btn'], callback_data="check_subscription")])
    toggle_text = "🇬🇧 English" if lang == 'fa' else "🇮🇷 فارسی"
    keyboard.append([InlineKeyboardButton(toggle_text, callback_data="toggle_lang")])
    return InlineKeyboardMarkup(keyboard)

# ========== CONVERSION FUNCTIONS ==========
async def convert_video_to_mp4(input_path: str, output_path: str, quality: str):
    try:
        height = int(quality.replace('p', ''))
        cmd = ['ffmpeg', '-i', input_path, '-vf', f'scale=-2:{height}', '-c:v', 'libx264', '-preset', 'fast', '-crf', '22', '-c:a', 'aac', '-b:a', '128k', '-movflags', '+faststart', '-y', output_path]
        result = subprocess.run(cmd, capture_output=True, text=True)
        return result.returncode == 0
    except Exception as e:
        logger.error(f"Video conversion error: {e}")
        return False

async def convert_image(input_path: str, output_path: str, output_format: str):
    try:
        img = Image.open(input_path)
        if output_format.lower() == 'jpg':
            img = img.convert('RGB')
            img.save(output_path, 'JPEG', quality=95)
        elif output_format.lower() == 'png':
            img.save(output_path, 'PNG')
        return True
    except Exception as e:
        logger.error(f"Image conversion error: {e}")
        return False

async def download_link_to_file(url: str, output_path: str):
    try:
        response = requests.get(url, stream=True, timeout=30)
        if response.status_code == 200:
            with open(output_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)
            return True
        return False
    except Exception as e:
        logger.error(f"Link download error: {e}")
        return False

def images_to_pdf(image_paths, output_path):
    """Convert multiple images to PDF"""
    try:
        images = []
        for path in image_paths:
            img = Image.open(path)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            images.append(img)
        if images:
            images[0].save(output_path, "PDF", save_all=True, append_images=images[1:])
            return True
        return False
    except Exception as e:
        logger.error(f"Images to PDF error: {e}")
        return False

def pdf_to_images(pdf_path, output_dir):
    """Convert PDF to images"""
    try:
        doc = fitz.open(pdf_path)
        image_paths = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            pix = page.get_pixmap()
            image_path = os.path.join(output_dir, f"page_{page_num + 1}.png")
            pix.save(image_path)
            image_paths.append(image_path)
        doc.close()
        return image_paths
    except Exception as e:
        logger.error(f"PDF to images error: {e}")
        return None

def generate_qr_code(text, output_path):
    """Generate QR code from text"""
    try:
        qr = qrcode.QRCode(version=1, box_size=10, border=5)
        qr.add_data(text)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        img.save(output_path)
        return True
    except Exception as e:
        logger.error(f"QR generate error: {e}")
        return False

def read_qr_code(image_path):
    """Read QR code from image"""
    try:
        img = Image.open(image_path)
        decoded = decode(img)
        if decoded:
            return decoded[0].data.decode('utf-8')
        return None
    except Exception as e:
        logger.error(f"QR read error: {e}")
        return None

def create_zip(file_paths, output_path):
    """Create ZIP archive from files"""
    try:
        with zipfile.ZipFile(output_path, 'w') as zipf:
            for file_path in file_paths:
                zipf.write(file_path, os.path.basename(file_path))
        return True
    except Exception as e:
        logger.error(f"ZIP create error: {e}")
        return False

def extract_zip(zip_path, output_dir):
    """Extract ZIP archive"""
    try:
        with zipfile.ZipFile(zip_path, 'r') as zipf:
            zipf.extractall(output_dir)
        return True
    except Exception as e:
        logger.error(f"ZIP extract error: {e}")
        return False

def doc_to_pdf(input_path, output_path):
    """Convert Word/PPT/Excel to PDF"""
    try:
        ext = os.path.splitext(input_path)[1].lower()
        if ext == '.docx':
            doc = Document(input_path)
            c = canvas.Canvas(output_path, pagesize=letter)
            # Simple conversion - extracts text
            for para in doc.paragraphs:
                if para.text:
                    c.drawString(50, 750, para.text[:100])
            c.save()
            return True
        elif ext == '.pptx':
            prs = Presentation(input_path)
            c = canvas.Canvas(output_path, pagesize=letter)
            y = 750
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        c.drawString(50, y, shape.text[:100])
                        y -= 20
                y -= 50
            c.save()
            return True
        elif ext in ['.xlsx', '.xls']:
            wb = load_workbook(input_path)
            c = canvas.Canvas(output_path, pagesize=letter)
            y = 750
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell else "" for cell in row)
                    if row_text:
                        c.drawString(50, y, row_text[:100])
                        y -= 20
                y -= 50
            c.save()
            return True
        return False
    except Exception as e:
        logger.error(f"Document to PDF error: {e}")
        return False

def text_to_pdf(text, output_path):
    """Convert text to PDF"""
    try:
        c = canvas.Canvas(output_path, pagesize=letter)
        y = 750
        for line in text.split('\n'):
            if y < 50:
                c.showPage()
                y = 750
            c.drawString(50, y, line[:100])
            y -= 20
        c.save()
        return True
    except Exception as e:
        logger.error(f"Text to PDF error: {e}")
        return False

def extract_audio_from_video(input_path, output_path):
    """Extract audio from video using FFmpeg"""
    try:
        cmd = ['ffmpeg', '-i', input_path, '-vn', '-acodec', 'mp3', '-ab', '192k', '-y', output_path]
        result = subprocess.run(cmd, capture_output=True, text=True)
        return result.returncode == 0
    except Exception as e:
        logger.error(f"Audio extraction error: {e}")
        return False

def unit_converter(value, from_unit, to_unit, category):
    """Convert units"""
    conversions = {
        'length': {
            'meter': 1, 'kilometer': 1000, 'centimeter': 0.01,
            'millimeter': 0.001, 'mile': 1609.34, 'yard': 0.9144,
            'foot': 0.3048, 'inch': 0.0254
        },
        'weight': {
            'kilogram': 1, 'gram': 0.001, 'milligram': 0.000001,
            'pound': 0.453592, 'ounce': 0.0283495, 'ton': 1000
        },
        'temperature': {
            'celsius': 'c', 'fahrenheit': 'f', 'kelvin': 'k'
        },
        'area': {
            'square_meter': 1, 'square_kilometer': 1000000,
            'square_mile': 2589988, 'acre': 4046.86,
            'hectare': 10000, 'square_foot': 0.092903
        },
        'volume': {
            'liter': 1, 'milliliter': 0.001, 'gallon': 3.78541,
            'quart': 0.946353, 'pint': 0.473176, 'cup': 0.236588
        }
    }
    
    try:
        if category == 'temperature':
            if from_unit == 'celsius' and to_unit == 'fahrenheit':
                return (value * 9/5) + 32
            elif from_unit == 'celsius' and to_unit == 'kelvin':
                return value + 273.15
            elif from_unit == 'fahrenheit' and to_unit == 'celsius':
                return (value - 32) * 5/9
            elif from_unit == 'fahrenheit' and to_unit == 'kelvin':
                return (value - 32) * 5/9 + 273.15
            elif from_unit == 'kelvin' and to_unit == 'celsius':
                return value - 273.15
            elif from_unit == 'kelvin' and to_unit == 'fahrenheit':
                return (value - 273.15) * 9/5 + 32
            return value
        else:
            base_value = value * conversions[category][from_unit]
            result = base_value / conversions[category][to_unit]
            return result
    except Exception as e:
        logger.error(f"Unit conversion error: {e}")
        return None

# ========== START & SUBSCRIPTION ==========
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    track_user(user.id, user.username, user.first_name, user.last_name, 'start')
    user_id = user.id
    if user_id not in user_lang:
        user_lang[user_id] = 'fa'
    lang = user_lang[user_id]
    
    message = LANG[lang]['welcome']
    for channel in CHANNELS:
        message += LANG[lang]['join_channels'].format(name=channel['name'])
    message += "\n" + LANG[lang]['help']
    
    await update.message.reply_text(
        message,
        reply_markup=build_channel_keyboard(lang),
        parse_mode="Markdown"
    )

async def check_subscription_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if await check_subscription(user_id, context):
        user_subscribed[user_id] = True
        await query.edit_message_text(
            LANG[lang]['subscribed'],
            reply_markup=build_main_menu(lang),
            parse_mode="Markdown"
        )
    else:
        await query.edit_message_text(
            LANG[lang]['not_subscribed'],
            reply_markup=build_channel_keyboard(lang),
            parse_mode="Markdown"
        )

async def toggle_language(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    current = user_lang.get(user_id, 'fa')
    new_lang = 'en' if current == 'fa' else 'fa'
    user_lang[user_id] = new_lang
    
    # Check if user is subscribed
    if await check_subscription(user_id, context):
        await query.edit_message_text(
            LANG[new_lang]['subscribed'],
            reply_markup=build_main_menu(new_lang),
            parse_mode="Markdown"
        )
    else:
        message = LANG[new_lang]['welcome']
        for channel in CHANNELS:
            message += LANG[new_lang]['join_channels'].format(name=channel['name'])
        message += "\n" + LANG[new_lang]['help']
        await query.edit_message_text(
            message,
            reply_markup=build_channel_keyboard(new_lang),
            parse_mode="Markdown"
        )

# ========== MENU NAVIGATION ==========
async def main_menu(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    # Clear session
    if user_id in user_sessions:
        del user_sessions[user_id]
    
    await query.edit_message_text(
        LANG[lang]['main_menu'],
        reply_markup=build_main_menu(lang),
        parse_mode="Markdown"
    )

async def menu_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    data = query.data
    lang = user_lang.get(user_id, 'fa')
    
    # Check subscription first
    if not await check_subscription(user_id, context):
        await query.edit_message_text(
            LANG[lang]['not_allowed'],
            reply_markup=build_channel_keyboard(lang),
            parse_mode="Markdown"
        )
        return
    
    # Handle menu selections
    menu_actions = {
        'menu_video_convert': ('🎬', 'send_video', 'video_conversion', build_back_button),
        'menu_audio_convert': ('🎵', 'send_audio', 'audio_conversion', build_back_button),
        'menu_youtube': ('📥', 'send_youtube', 'youtube_download', build_back_button),
        'menu_instagram': ('📸', 'send_instagram', 'instagram_download', build_back_button),
        'menu_image_convert': ('🖼', 'send_image', 'image_conversion', build_back_button),
        'menu_image_to_pdf': ('📄', 'send_images', 'image_to_pdf', build_back_button),
        'menu_pdf_to_image': ('📄', 'send_pdf', 'pdf_to_image', build_back_button),
        'menu_qr_generate': ('🔲', 'send_qr_text', 'qr_generate', build_back_button),
        'menu_qr_read': ('🔍', 'send_qr_image', 'qr_read', build_back_button),
        'menu_zip': ('📦', 'send_zip_files', 'zip_compress', build_back_button),
        'menu_unzip': ('📂', 'send_unzip', 'zip_extract', build_back_button),
        'menu_doc_to_pdf': ('📄', 'send_document', 'doc_to_pdf', build_back_button),
        'menu_text_to_pdf': ('📝', 'send_text', 'text_to_pdf', build_back_button),
        'menu_audio_extract': ('🎬', 'send_video', 'audio_extract', build_back_button),
        'menu_link_download': ('🔗', 'send_link', 'link_download', build_back_button),
        'menu_unit_convert': ('📏', '', 'unit_convert', build_unit_converter_keyboard),
    }
    
    if data in menu_actions:
        icon, prompt, session_type, keyboard_func = menu_actions[data]
        
        # Set user session
        user_sessions[user_id] = {'state': session_type, 'files': []}
        
        if data == 'menu_unit_convert':
            await query.edit_message_text(
                f"{icon} **{LANG[lang]['menu_unit_convert']}**\n\n"
                "📏 Select conversion type:",
                reply_markup=keyboard_func(lang),
                parse_mode="Markdown"
            )
        elif data == 'menu_image_to_pdf':
            await query.edit_message_text(
                f"{icon} **{LANG[lang]['menu_image_to_pdf']}**\n\n"
                f"{LANG[lang]['send_images']}",
                reply_markup=keyboard_func(lang),
                parse_mode="Markdown"
            )
        elif data == 'menu_zip':
            await query.edit_message_text(
                f"{icon} **{LANG[lang]['menu_zip']}**\n\n"
                f"{LANG[lang]['send_zip_files']}",
                reply_markup=keyboard_func(lang),
                parse_mode="Markdown"
            )
        else:
            await query.edit_message_text(
                f"{icon} **{LANG[lang][menu_actions[data][0]]}**\n\n"
                f"{LANG[lang][prompt]}",
                reply_markup=keyboard_func(lang),
                parse_mode="Markdown"
            )

# ========== UNIT CONVERTER ==========
async def unit_convert_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    data = query.data
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if data == "main_menu":
        await main_menu(update, context)
        return
    
    # Store unit category in session
    if user_id not in user_sessions:
        user_sessions[user_id] = {}
    user_sessions[user_id]['unit_category'] = data.replace('unit_', '')
    user_sessions[user_id]['state'] = 'unit_value'
    
    await query.edit_message_text(
        "📏 **Unit Converter**\n\n"
        "Please enter the value and units in this format:\n"
        "`10 meter to kilometer`\n"
        "`25 celsius to fahrenheit`\n\n"
        "Available units:\n"
        "Length: meter, kilometer, centimeter, millimeter, mile, yard, foot, inch\n"
        "Weight: kilogram, gram, milligram, pound, ounce, ton\n"
        "Temperature: celsius, fahrenheit, kelvin\n"
        "Area: square_meter, square_kilometer, square_mile, acre, hectare, square_foot\n"
        "Volume: liter, milliliter, gallon, quart, pint, cup",
        reply_markup=build_back_button(lang),
        parse_mode="Markdown"
    )

# ========== FILE HANDLING ==========
async def handle_file(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if not await check_subscription(user_id, context):
        await update.message.reply_text(
            LANG[lang]['not_allowed'],
            reply_markup=build_channel_keyboard(lang),
            parse_mode="Markdown"
        )
        return
    
    session = user_sessions.get(user_id, {})
    state = session.get('state', '')
    
    # Handle different states
    if state == 'video_conversion':
        await handle_video_conversion(update, context)
    elif state == 'audio_conversion':
        await handle_audio_conversion(update, context)
    elif state == 'image_conversion':
        await handle_image_conversion(update, context)
    elif state == 'image_to_pdf':
        await handle_image_to_pdf(update, context)
    elif state == 'pdf_to_image':
        await handle_pdf_to_image(update, context)
    elif state == 'qr_read':
        await handle_qr_read(update, context)
    elif state == 'zip_extract':
        await handle_zip_extract(update, context)
    elif state == 'doc_to_pdf':
        await handle_doc_to_pdf(update, context)
    elif state == 'audio_extract':
        await handle_audio_extract(update, context)
    elif state == 'zip_compress':
        await handle_zip_compress(update, context)
    elif state == 'link_download':
        await handle_link(update, context, update.message.text)
    elif state == 'unit_value':
        await handle_unit_conversion(update, context)
    elif state == 'text_to_pdf':
        await handle_text_to_pdf(update, context)
    elif state == 'qr_generate':
        await handle_qr_generate(update, context)
    elif state == 'youtube_download':
        await handle_youtube(update, context)
    elif state == 'instagram_download':
        await handle_instagram(update, context)
    else:
        await update.message.reply_text(
            "🔙 Please select an option from the menu first.",
            reply_markup=build_main_menu(lang)
        )

# ========== FEATURE HANDLERS ==========
async def handle_video_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    video = update.message.video
    if not video:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_video'])
        return
    
    file_id = video.file_id
    ref_id = str(uuid.uuid4())[:8]
    temp_storage[ref_id] = file_id
    
    keyboard = [
        [InlineKeyboardButton("720p", callback_data=f"vidconv_720_{ref_id}")],
        [InlineKeyboardButton("1080p", callback_data=f"vidconv_1080_{ref_id}")],
        [InlineKeyboardButton("480p", callback_data=f"vidconv_480_{ref_id}")],
        [InlineKeyboardButton(LANG[lang]['back'], callback_data="main_menu")]
    ]
    await update.message.reply_text(
        LANG[lang]['video_conversion'],
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode="Markdown"
    )

async def handle_audio_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    file = update.message.audio or update.message.voice
    if not file:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_audio'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        file_obj = await context.bot.get_file(file.file_id)
        with tempfile.NamedTemporaryFile(delete=False) as input_file:
            input_path = input_file.name
            await file_obj.download_to_drive(input_path)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as output_file:
            output_path = output_file.name
        
        audio = AudioSegment.from_file(input_path)
        audio.export(output_path, format="mp3", bitrate="192k")
        
        with open(output_path, 'rb') as mp3_file:
            await update.message.reply_audio(
                audio=mp3_file,
                filename="converted.mp3",
                performer="Bot",
                title="Converted Audio"
            )
        await processing_msg.delete()
        os.unlink(input_path)
        os.unlink(output_path)
    except Exception as e:
        logger.error(f"Audio conversion error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_image_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if update.message.photo:
        file_id = update.message.photo[-1].file_id
        current_format = 'jpg'
    else:
        doc = update.message.document
        if not doc or not doc.mime_type or not doc.mime_type.startswith('image/'):
            await update.message.reply_text("⚠️ " + LANG[lang]['send_image'])
            return
        file_id = doc.file_id
        current_format = doc.mime_type.split('/')[-1]
    
    ref_id = str(uuid.uuid4())[:8]
    temp_storage[ref_id] = file_id
    
    keyboard = []
    if current_format.lower() in ['jpg', 'jpeg']:
        keyboard.append([InlineKeyboardButton("PNG", callback_data=f"imgconv_png_{ref_id}")])
    elif current_format.lower() == 'png':
        keyboard.append([InlineKeyboardButton("JPG", callback_data=f"imgconv_jpg_{ref_id}")])
    else:
        keyboard.append([InlineKeyboardButton("JPG", callback_data=f"imgconv_jpg_{ref_id}")])
        keyboard.append([InlineKeyboardButton("PNG", callback_data=f"imgconv_png_{ref_id}")])
    keyboard.append([InlineKeyboardButton(LANG[lang]['back'], callback_data="main_menu")])
    
    await update.message.reply_text(
        LANG[lang]['image_conversion'],
        reply_markup=InlineKeyboardMarkup(keyboard),
        parse_mode="Markdown"
    )

async def handle_image_to_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if user_id not in user_sessions:
        user_sessions[user_id] = {'state': 'image_to_pdf', 'files': []}
    
    session = user_sessions[user_id]
    
    if update.message.photo:
        file_id = update.message.photo[-1].file_id
        session['files'].append(file_id)
    elif update.message.document:
        doc = update.message.document
        if doc.mime_type and doc.mime_type.startswith('image/'):
            session['files'].append(doc.file_id)
        else:
            await update.message.reply_text("⚠️ " + LANG[lang]['send_image'])
            return
    else:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_image'])
        return
    
    if len(session['files']) >= 5:
        # Create PDF
        await create_pdf_from_images(update, context)
    else:
        keyboard = [
            [InlineKeyboardButton("📄 Create PDF", callback_data="create_pdf")],
            [InlineKeyboardButton(LANG[lang]['back'], callback_data="main_menu")]
        ]
        await update.message.reply_text(
            f"✅ Image added! ({len(session['files'])}/5)\n\n{LANG[lang]['send_images']}",
            reply_markup=InlineKeyboardMarkup(keyboard)
        )

async def create_pdf_from_images(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    session = user_sessions.get(user_id, {})
    files = session.get('files', [])
    
    if not files:
        await query.edit_message_text("❌ No images found.")
        return
    
    await query.edit_message_text(LANG[lang]['processing'])
    
    try:
        image_paths = []
        for file_id in files:
            file = await context.bot.get_file(file_id)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp:
                path = tmp.name
                await file.download_to_drive(path)
                image_paths.append(path)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            output_path = tmp.name
        
        success = images_to_pdf(image_paths, output_path)
        
        if success:
            with open(output_path, 'rb') as pdf_file:
                await query.message.reply_document(
                    document=pdf_file,
                    filename="images.pdf",
                    caption=LANG[lang]['pdf_created']
                )
            await query.edit_message_text(LANG[lang]['download_complete'])
        else:
            await query.edit_message_text(LANG[lang]['error'])
        
        # Clean up
        for path in image_paths:
            if os.path.exists(path):
                os.unlink(path)
        if os.path.exists(output_path):
            os.unlink(output_path)
        
        # Clear session
        user_sessions[user_id] = {'state': None, 'files': []}
        
    except Exception as e:
        logger.error(f"PDF creation error: {e}")
        await query.edit_message_text(LANG[lang]['error'])

async def handle_pdf_to_image(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    doc = update.message.document
    if not doc or not doc.mime_type or doc.mime_type != 'application/pdf':
        await update.message.reply_text("⚠️ " + LANG[lang]['send_pdf'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        file_obj = await context.bot.get_file(doc.file_id)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as pdf_file:
            pdf_path = pdf_file.name
            await file_obj.download_to_drive(pdf_path)
        
        output_dir = tempfile.mkdtemp()
        image_paths = pdf_to_images(pdf_path, output_dir)
        
        if image_paths:
            for img_path in image_paths:
                with open(img_path, 'rb') as img:
                    await update.message.reply_photo(photo=img)
            await processing_msg.delete()
            shutil.rmtree(output_dir)
        else:
            await processing_msg.edit_text(LANG[lang]['error'])
        
        os.unlink(pdf_path)
    except Exception as e:
        logger.error(f"PDF to image error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_qr_generate(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    text = update.message.text
    if not text:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_qr_text'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            output_path = tmp.name
        
        success = generate_qr_code(text, output_path)
        if success:
            with open(output_path, 'rb') as qr_file:
                await update.message.reply_photo(
                    photo=qr_file,
                    caption=LANG[lang]['qr_created']
                )
            await processing_msg.delete()
        else:
            await processing_msg.edit_text(LANG[lang]['error'])
        os.unlink(output_path)
    except Exception as e:
        logger.error(f"QR generate error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_qr_read(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if update.message.photo:
        file_id = update.message.photo[-1].file_id
    elif update.message.document:
        doc = update.message.document
        if doc.mime_type and doc.mime_type.startswith('image/'):
            file_id = doc.file_id
        else:
            await update.message.reply_text("⚠️ " + LANG[lang]['send_qr_image'])
            return
    else:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_qr_image'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        file_obj = await context.bot.get_file(file_id)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp:
            image_path = tmp.name
            await file_obj.download_to_drive(image_path)
        
        result = read_qr_code(image_path)
        if result:
            await update.message.reply_text(
                LANG[lang]['qr_read_result'].format(text=result)
            )
        else:
            await update.message.reply_text("❌ No QR Code found in image.")
        await processing_msg.delete()
        os.unlink(image_path)
    except Exception as e:
        logger.error(f"QR read error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_zip_compress(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if user_id not in user_sessions:
        user_sessions[user_id] = {'state': 'zip_compress', 'files': []}
    
    session = user_sessions[user_id]
    
    doc = update.message.document
    if doc:
        session['files'].append(doc.file_id)
    else:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_zip_files'])
        return
    
    if len(session['files']) >= 5:
        await create_zip_from_files(update, context)
    else:
        keyboard = [
            [InlineKeyboardButton("📦 Compress", callback_data="create_zip")],
            [InlineKeyboardButton(LANG[lang]['back'], callback_data="main_menu")]
        ]
        await update.message.reply_text(
            f"✅ File added! ({len(session['files'])}/5)\n\n{LANG[lang]['send_zip_files']}",
            reply_markup=InlineKeyboardMarkup(keyboard)
        )

async def create_zip_from_files(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    session = user_sessions.get(user_id, {})
    files = session.get('files', [])
    
    if not files:
        await query.edit_message_text("❌ No files found.")
        return
    
    await query.edit_message_text(LANG[lang]['processing'])
    
    try:
        file_paths = []
        for file_id in files:
            file = await context.bot.get_file(file_id)
            with tempfile.NamedTemporaryFile(delete=False) as tmp:
                path = tmp.name
                await file.download_to_drive(path)
                file_paths.append(path)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".zip") as tmp:
            output_path = tmp.name
        
        success = create_zip(file_paths, output_path)
        
        if success:
            with open(output_path, 'rb') as zip_file:
                await query.message.reply_document(
                    document=zip_file,
                    filename="archive.zip",
                    caption=LANG[lang]['zip_created']
                )
            await query.edit_message_text(LANG[lang]['download_complete'])
        else:
            await query.edit_message_text(LANG[lang]['error'])
        
        for path in file_paths:
            if os.path.exists(path):
                os.unlink(path)
        if os.path.exists(output_path):
            os.unlink(output_path)
        
        user_sessions[user_id] = {'state': None, 'files': []}
        
    except Exception as e:
        logger.error(f"ZIP create error: {e}")
        await query.edit_message_text(LANG[lang]['error'])

async def handle_zip_extract(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    doc = update.message.document
    if not doc or not doc.mime_type or doc.mime_type != 'application/zip':
        await update.message.reply_text("⚠️ " + LANG[lang]['send_unzip'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        file_obj = await context.bot.get_file(doc.file_id)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".zip") as zip_file:
            zip_path = zip_file.name
            await file_obj.download_to_drive(zip_path)
        
        output_dir = tempfile.mkdtemp()
        success = extract_zip(zip_path, output_dir)
        
        if success:
            # Send extracted files
            for filename in os.listdir(output_dir):
                file_path = os.path.join(output_dir, filename)
                if os.path.isfile(file_path):
                    with open(file_path, 'rb') as f:
                        await update.message.reply_document(
                            document=f,
                            filename=filename
                        )
            await processing_msg.delete()
            shutil.rmtree(output_dir)
        else:
            await processing_msg.edit_text(LANG[lang]['error'])
        
        os.unlink(zip_path)
    except Exception as e:
        logger.error(f"ZIP extract error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_doc_to_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    doc = update.message.document
    if not doc:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_document'])
        return
    
    ext = os.path.splitext(doc.file_name)[1].lower()
    if ext not in ['.docx', '.pptx', '.xlsx', '.xls']:
        await update.message.reply_text("⚠️ Please send a Word (.docx), PowerPoint (.pptx), or Excel (.xlsx) file.")
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        file_obj = await context.bot.get_file(doc.file_id)
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as input_file:
            input_path = input_file.name
            await file_obj.download_to_drive(input_path)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as output_file:
            output_path = output_file.name
        
        success = doc_to_pdf(input_path, output_path)
        
        if success:
            with open(output_path, 'rb') as pdf_file:
                await update.message.reply_document(
                    document=pdf_file,
                    filename=os.path.splitext(doc.file_name)[0] + ".pdf",
                    caption="📄 PDF created from document"
                )
            await processing_msg.delete()
        else:
            await processing_msg.edit_text(LANG[lang]['error'])
        
        os.unlink(input_path)
        if os.path.exists(output_path):
            os.unlink(output_path)
    except Exception as e:
        logger.error(f"Document to PDF error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_text_to_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    text = update.message.text
    if not text:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_text'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as output_file:
            output_path = output_file.name
        
        success = text_to_pdf(text, output_path)
        
        if success:
            with open(output_path, 'rb') as pdf_file:
                await update.message.reply_document(
                    document=pdf_file,
                    filename="text.pdf",
                    caption=LANG[lang]['pdf_created']
                )
            await processing_msg.delete()
        else:
            await processing_msg.edit_text(LANG[lang]['error'])
        
        if os.path.exists(output_path):
            os.unlink(output_path)
    except Exception as e:
        logger.error(f"Text to PDF error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_audio_extract(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    video = update.message.video
    if not video:
        await update.message.reply_text("⚠️ " + LANG[lang]['send_video'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        file_obj = await context.bot.get_file(video.file_id)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as input_file:
            input_path = input_file.name
            await file_obj.download_to_drive(input_path)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as output_file:
            output_path = output_file.name
        
        success = extract_audio_from_video(input_path, output_path)
        
        if success:
            with open(output_path, 'rb') as audio_file:
                await update.message.reply_audio(
                    audio=audio_file,
                    filename="extracted_audio.mp3",
                    performer="Bot",
                    title="Extracted Audio"
                )
            await processing_msg.delete()
        else:
            await processing_msg.edit_text(LANG[lang]['error'])
        
        os.unlink(input_path)
        if os.path.exists(output_path):
            os.unlink(output_path)
    except Exception as e:
        logger.error(f"Audio extract error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_unit_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    text = update.message.text
    if not text:
        await update.message.reply_text(
            "📏 Please enter conversion in format:\n`10 meter to kilometer`",
            parse_mode="Markdown"
        )
        return
    
    try:
        # Parse input: "10 meter to kilometer"
        parts = text.lower().split(' to ')
        if len(parts) != 2:
            await update.message.reply_text(
                "❌ Invalid format. Use:\n`10 meter to kilometer`",
                parse_mode="Markdown"
            )
            return
        
        value_str = parts[0].strip()
        to_unit = parts[1].strip()
        
        # Extract value and from_unit
        value_parts = value_str.split()
        if len(value_parts) < 2:
            await update.message.reply_text(
                "❌ Invalid format. Use:\n`10 meter to kilometer`",
                parse_mode="Markdown"
            )
            return
        
        value = float(value_parts[0])
        from_unit = ' '.join(value_parts[1:])
        
        # Determine category
        category = user_sessions.get(user_id, {}).get('unit_category', 'length')
        
        result = unit_converter(value, from_unit, to_unit, category)
        
        if result is not None:
            await update.message.reply_text(
                LANG[lang]['unit_result'].format(
                    result=f"{value} {from_unit} = {result:.4f} {to_unit}"
                ),
                reply_markup=build_back_button(lang)
            )
        else:
            await update.message.reply_text(
                "❌ Conversion failed. Please check your units.",
                reply_markup=build_back_button(lang)
            )
    except ValueError:
        await update.message.reply_text(
            "❌ Invalid number format. Use:\n`10 meter to kilometer`",
            parse_mode="Markdown"
        )
    except Exception as e:
        logger.error(f"Unit conversion error: {e}")
        await update.message.reply_text(LANG[lang]['error'])

async def handle_link(update: Update, context: ContextTypes.DEFAULT_TYPE, url: str):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if not url.startswith('http'):
        await update.message.reply_text("⚠️ " + LANG[lang]['send_link'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    try:
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            output_path = tmp.name
        
        success = await download_link_to_file(url, output_path)
        if success:
            file_size = os.path.getsize(output_path)
            if file_size > 50 * 1024 * 1024:
                await processing_msg.edit_text("❌ File too large for Telegram (max 50 MB).")
                os.unlink(output_path)
                return
            with open(output_path, 'rb') as f:
                await update.message.reply_document(
                    document=f,
                    filename=os.path.basename(url.split('/')[-1]) or 'file.bin',
                    caption="📁 File downloaded from link"
                )
            await processing_msg.delete()
        else:
            await processing_msg.edit_text(LANG[lang]['error'])
        if os.path.exists(output_path):
            os.unlink(output_path)
    except Exception as e:
        logger.error(f"Link download error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

# ========== YOUTUBE & INSTAGRAM ==========
async def handle_youtube(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    url = update.message.text.strip()
    
    if "youtube.com" not in url.lower() and "youtu.be" not in url.lower():
        await update.message.reply_text("⚠️ " + LANG[lang]['send_youtube'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    
    ydl_opts = {
        'quiet': True,
        'no_warnings': True,
        'extractor_args': {'youtube': {'player_client': ['tv', 'web'], 'skip': ['dash', 'hls']}},
        'headers': {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'},
        'ignoreerrors': True,
    }
    
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            if info is None:
                raise Exception("No info")
            title = info.get('title', 'Video')
            formats = info.get('formats', [])
            
            keyboard = []
            for res in ['1080p', '720p', '480p', '360p']:
                height = int(res.replace('p', ''))
                if any(f.get('height') == height for f in formats):
                    keyboard.append([InlineKeyboardButton(f"📹 {res}", callback_data=f"yt_video_{res}_{url}")])
            keyboard.append([InlineKeyboardButton("🎵 MP3", callback_data=f"yt_audio_{url}")])
            keyboard.append([InlineKeyboardButton(LANG[lang]['back'], callback_data="main_menu")])
            
            await processing_msg.edit_text(
                f"🎬 **{title}**\n\n{LANG[lang]['choose_quality']}",
                reply_markup=InlineKeyboardMarkup(keyboard),
                parse_mode="Markdown"
            )
    except Exception as e:
        logger.error(f"YouTube error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def handle_instagram(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    url = update.message.text.strip()
    
    if "instagram.com" not in url.lower():
        await update.message.reply_text("⚠️ " + LANG[lang]['send_instagram'])
        return
    
    processing_msg = await update.message.reply_text(LANG[lang]['processing'])
    
    ydl_opts = {
        'quiet': True,
        'no_warnings': True,
        'headers': {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'},
        'ignoreerrors': True,
    }
    
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            if info is None:
                raise Exception("No info")
            title = info.get('title', 'Instagram Reel')
            title = title.replace('Instagram', '').strip()
            
            keyboard = [
                [InlineKeyboardButton("📹 Video (MP4)", callback_data=f"insta_video_{url}")],
                [InlineKeyboardButton("🎵 Audio (MP3)", callback_data=f"insta_audio_{url}")],
                [InlineKeyboardButton(LANG[lang]['back'], callback_data="main_menu")]
            ]
            await processing_msg.edit_text(
                f"📸 **{title}**\n\nChoose option:",
                reply_markup=InlineKeyboardMarkup(keyboard),
                parse_mode="Markdown"
            )
    except Exception as e:
        logger.error(f"Instagram error: {e}")
        await processing_msg.edit_text(LANG[lang]['error'])

async def youtube_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    data = query.data
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    parts = data.split('_')
    if parts[0] == "yt":
        if parts[1] == "video":
            resolution = parts[2]
            url = '_'.join(parts[3:])
            await download_youtube_video(query, url, resolution, lang)
        elif parts[1] == "audio":
            url = '_'.join(parts[2:])
            await download_youtube_audio(query, url, lang)

async def download_youtube_video(query, url, resolution, lang):
    await query.edit_message_text(LANG[lang]['processing'])
    try:
        height = int(resolution.replace('p', ''))
        ydl_opts = {
            'format': f'bestvideo[height<={height}][ext=mp4]+bestaudio[ext=m4a]/best[height<={height}][ext=mp4]',
            'outtmpl': '%(title)s.%(ext)s',
            'quiet': True,
            'no_warnings': True,
            'extractor_args': {'youtube': {'player_client': ['tv', 'web'], 'skip': ['dash', 'hls']}},
            'headers': {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'},
            'ignoreerrors': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            if info is None:
                raise Exception("Download failed")
            filepath = ydl.prepare_filename(info)
            with open(filepath, 'rb') as video_file:
                await query.message.reply_video(
                    video=video_file,
                    caption=f"📹 {info.get('title', 'Video')}\nQuality: {resolution}"
                )
            os.unlink(filepath)
        await query.edit_message_text(LANG[lang]['download_complete'])
    except Exception as e:
        logger.error(f"YouTube download error: {e}")
        await query.edit_message_text(LANG[lang]['error'])

async def download_youtube_audio(query, url, lang):
    await query.edit_message_text(LANG[lang]['processing'])
    try:
        ydl_opts = {
            'format': 'bestaudio[ext=m4a]/bestaudio/best',
            'postprocessors': [{'key': 'FFmpegExtractAudio', 'preferredcodec': 'mp3', 'preferredquality': '192'}],
            'outtmpl': '%(title)s.%(ext)s',
            'quiet': True,
            'no_warnings': True,
            'extractor_args': {'youtube': {'player_client': ['tv', 'web'], 'skip': ['dash', 'hls']}},
            'headers': {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'},
            'ignoreerrors': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            if info is None:
                raise Exception("Download failed")
            filepath = ydl.prepare_filename(info).replace('.webm', '.mp3')
            with open(filepath, 'rb') as audio_file:
                await query.message.reply_audio(
                    audio=audio_file,
                    performer="YouTube",
                    title=info.get('title', 'Audio'),
                    caption="🎵 Extracted from YouTube"
                )
            os.unlink(filepath)
        await query.edit_message_text(LANG[lang]['download_complete'])
    except Exception as e:
        logger.error(f"YouTube audio download error: {e}")
        await query.edit_message_text(LANG[lang]['error'])

async def instagram_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    data = query.data
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    parts = data.split('_')
    if parts[0] == "insta":
        if parts[1] == "video":
            url = '_'.join(parts[2:])
            await download_instagram_video(query, url, lang)
        elif parts[1] == "audio":
            url = '_'.join(parts[2:])
            await download_instagram_audio(query, url, lang)

async def download_instagram_video(query, url, lang):
    await query.edit_message_text(LANG[lang]['processing'])
    try:
        ydl_opts = {
            'format': 'best[ext=mp4]/best',
            'outtmpl': '%(title)s.%(ext)s',
            'quiet': True,
            'no_warnings': True,
            'headers': {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'},
            'ignoreerrors': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            if info is None:
                raise Exception("Download failed")
            filepath = ydl.prepare_filename(info)
            title = info.get('title', 'Instagram Reel')
            title = title.replace('Instagram', '').strip()
            with open(filepath, 'rb') as video_file:
                await query.message.reply_video(
                    video=video_file,
                    caption=f"📸 {title}"
                )
            os.unlink(filepath)
        await query.edit_message_text(LANG[lang]['download_complete'])
    except Exception as e:
        logger.error(f"Instagram download error: {e}")
        await query.edit_message_text(LANG[lang]['error'])

async def download_instagram_audio(query, url, lang):
    await query.edit_message_text(LANG[lang]['processing'])
    try:
        ydl_opts = {
            'format': 'bestaudio[ext=m4a]/bestaudio/best',
            'postprocessors': [{'key': 'FFmpegExtractAudio', 'preferredcodec': 'mp3', 'preferredquality': '192'}],
            'outtmpl': '%(title)s.%(ext)s',
            'quiet': True,
            'no_warnings': True,
            'headers': {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'},
            'ignoreerrors': True,
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            if info is None:
                raise Exception("Download failed")
            filepath = ydl.prepare_filename(info).replace('.webm', '.mp3')
            title = info.get('title', 'Instagram Reel')
            title = title.replace('Instagram', '').strip()
            with open(filepath, 'rb') as audio_file:
                await query.message.reply_audio(
                    audio=audio_file,
                    performer="Instagram",
                    title=title,
                    caption="🎵 Audio from Instagram"
                )
            os.unlink(filepath)
        await query.edit_message_text(LANG[lang]['download_complete'])
    except Exception as e:
        logger.error(f"Instagram audio download error: {e}")
        await query.edit_message_text(LANG[lang]['error'])

# ========== VIDEO CONVERSION CALLBACK ==========
async def video_conversion_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    data = query.data
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if data == "main_menu":
        await main_menu(update, context)
        return
    
    parts = data.split('_')
    if parts[0] == "vidconv":
        quality = parts[1]
        ref_id = parts[2]
        file_id = temp_storage.get(ref_id)
        if not file_id:
            await query.edit_message_text(LANG[lang]['file_not_found'])
            return
        
        await query.edit_message_text(LANG[lang]['processing'])
        
        try:
            file = await context.bot.get_file(file_id)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as input_file:
                input_path = input_file.name
                await file.download_to_drive(input_path)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as output_file:
                output_path = output_file.name
            
            success = await convert_video_to_mp4(input_path, output_path, quality)
            if success:
                with open(output_path, 'rb') as vid_file:
                    await query.message.reply_video(
                        video=vid_file,
                        caption=f"📹 MP4 {quality}",
                        supports_streaming=True
                    )
                await query.edit_message_text(LANG[lang]['download_complete'])
            else:
                await query.edit_message_text(LANG[lang]['error'])
            os.unlink(input_path)
            if os.path.exists(output_path):
                os.unlink(output_path)
            temp_storage.pop(ref_id, None)
        except Exception as e:
            logger.error(f"Video conversion error: {e}")
            await query.edit_message_text(LANG[lang]['error'])

# ========== IMAGE CONVERSION CALLBACK ==========
async def image_conversion_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    data = query.data
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    if data == "main_menu":
        await main_menu(update, context)
        return
    
    parts = data.split('_')
    if parts[0] == "imgconv":
        output_format = parts[1]
        ref_id = parts[2]
        file_id = temp_storage.get(ref_id)
        if not file_id:
            await query.edit_message_text(LANG[lang]['file_not_found'])
            return
        
        await query.edit_message_text(LANG[lang]['processing'])
        
        try:
            file = await context.bot.get_file(file_id)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".img") as input_file:
                input_path = input_file.name
                await file.download_to_drive(input_path)
            ext = '.' + output_format
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as output_file:
                output_path = output_file.name
            
            success = await convert_image(input_path, output_path, output_format)
            if success:
                with open(output_path, 'rb') as img_file:
                    if output_format == 'jpg':
                        await query.message.reply_photo(photo=img_file, caption="🖼 JPG Image")
                    else:
                        await query.message.reply_document(document=img_file, filename=f"image.{output_format}", caption="🖼 PNG Image")
                await query.edit_message_text(LANG[lang]['download_complete'])
            else:
                await query.edit_message_text(LANG[lang]['error'])
            os.unlink(input_path)
            if os.path.exists(output_path):
                os.unlink(output_path)
            temp_storage.pop(ref_id, None)
        except Exception as e:
            logger.error(f"Image conversion error: {e}")
            await query.edit_message_text(LANG[lang]['error'])

# ========== CALLBACK HANDLER ==========
async def callback_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    data = query.data
    user_id = query.from_user.id
    lang = user_lang.get(user_id, 'fa')
    
    # Handle main menu navigation
    if data == "main_menu":
        await main_menu(update, context)
        return
    
    # Handle subscription check
    if data == "check_subscription":
        await check_subscription_callback(update, context)
        return
    
    # Handle language toggle
    if data == "toggle_lang":
        await toggle_language(update, context)
        return
    
    # Handle create_pdf
    if data == "create_pdf":
        await create_pdf_from_images(update, context)
        return
    
    # Handle create_zip
    if data == "create_zip":
        await create_zip_from_files(update, context)
        return
    
    # Handle unit converter
    if data.startswith("unit_"):
        await unit_convert_callback(update, context)
        return
    
    # Handle video conversion
    if data.startswith("vidconv_"):
        await video_conversion_callback(update, context)
        return
    
    # Handle image conversion
    if data.startswith("imgconv_"):
        await image_conversion_callback(update, context)
        return
    
    # Handle YouTube
    if data.startswith("yt_"):
        await youtube_callback(update, context)
        return
    
    # Handle Instagram
    if data.startswith("insta_"):
        await instagram_callback(update, context)
        return
    
    # Handle menu selections
    menu_options = [
        'menu_video_convert', 'menu_audio_convert', 'menu_youtube', 'menu_instagram',
        'menu_image_convert', 'menu_image_to_pdf', 'menu_pdf_to_image', 'menu_qr_generate',
        'menu_qr_read', 'menu_zip', 'menu_unzip', 'menu_doc_to_pdf', 'menu_text_to_pdf',
        'menu_audio_extract', 'menu_unit_convert', 'menu_link_download'
    ]
    
    if data in menu_options:
        await menu_handler(update, context)

# ========== MAIN MESSAGE HANDLER ==========
async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    lang = user_lang.get(user_id, 'fa')
    text = update.message.text or ""
    
    # Check if user is in a session
    session = user_sessions.get(user_id, {})
    state = session.get('state', '')
    
    if text and text.startswith('/'):
        return
    
    # Check for direct link
    if text.startswith('http') and not state:
        await handle_link(update, context, text)
        return
    
    # Handle based on session state
    if state in ['video_conversion', 'audio_conversion', 'image_conversion', 'image_to_pdf',
                 'pdf_to_image', 'qr_read', 'zip_extract', 'doc_to_pdf', 'audio_extract',
                 'zip_compress', 'link_download', 'unit_value', 'text_to_pdf', 'qr_generate']:
        await handle_file(update, context)
    elif state in ['youtube_download']:
        await handle_youtube(update, context)
    elif state in ['instagram_download']:
        await handle_instagram(update, context)
    else:
        # Check if user is subscribed before showing menu
        if await check_subscription(user_id, context):
            await update.message.reply_text(
                LANG[lang]['main_menu'],
                reply_markup=build_main_menu(lang),
                parse_mode="Markdown"
            )
        else:
            await update.message.reply_text(
                LANG[lang]['not_allowed'],
                reply_markup=build_channel_keyboard(lang),
                parse_mode="Markdown"
            )

# ========== STATS (Admin) ==========
async def stats(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if user_id not in ADMIN_IDS:
        await update.message.reply_text("❌ You are not authorized to view stats.")
        return
    
    total_users = len(user_data)
    total_messages = sum(u['total_messages'] for u in user_data.values())
    
    stats_text = f"""
📊 **Bot Statistics**

👥 Total Users: {total_users}
💬 Total Messages: {total_messages}

**Actions Summary:**
"""
    action_counts = {}
    for user in user_data.values():
        for action, count in user.get('actions', {}).items():
            action_counts[action] = action_counts.get(action, 0) + count
    
    for action, count in action_counts.items():
        stats_text += f"• {action}: {count}\n"
    
    stats_text += "\n**Top 5 Most Active Users:**\n"
    sorted_users = sorted(user_data.values(), key=lambda x: x['total_messages'], reverse=True)[:5]
    for i, user in enumerate(sorted_users, 1):
        name = user.get('first_name', 'Unknown')
        username = user.get('username', 'No username')
        stats_text += f"{i}. {name} (@{username}) - {user['total_messages']} messages\n"
    
    await update.message.reply_text(stats_text, parse_mode="Markdown")

# ========== MAIN ==========
def main():
    app = Application.builder().token(TOKEN).build()
    
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("stats", stats))
    app.add_handler(CallbackQueryHandler(callback_handler))
    app.add_handler(MessageHandler(filters.ALL & ~filters.COMMAND, handle_message))
    
    print("🤖 Multipurpose Bot is running...")
    print("✅ All features loaded successfully!")
    print("📊 17+ Features available:")
    print("  🎬 Video Converter (720p/1080p/480p)")
    print("  🎵 Audio Converter (to MP3)")
    print("  📥 YouTube Downloader")
    print("  📸 Instagram Downloader")
    print("  🖼 Image Converter (JPG/PNG)")
    print("  📄 Image to PDF")
    print("  📄 PDF to Image")
    print("  🔲 QR Code Generator")
    print("  🔍 QR Code Reader")
    print("  📦 ZIP Compressor")
    print("  📂 ZIP Extractor")
    print("  📄 Document to PDF (Word/PPT/Excel)")
    print("  📝 Text to PDF")
    print("  🎬 Extract Audio from Video")
    print("  📏 Unit Converter")
    print("  🔗 Link Downloader")
    print("  🌐 Language Toggle (Farsi/English)")
    
    app.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == "__main__":
    main()
